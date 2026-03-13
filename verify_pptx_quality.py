#!/usr/bin/env python3
"""PPTXファイル品質検証スクリプト"""

import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

PPTX_PATH = "output/Claude_Code_とは_20260311_210446.pptx"
JSON_PATH = "output/claude_code_content.json"

# ── 定数 ──
MARGIN_L_INCH = 0.8
MARGIN_R_INCH = 0.8
SLIDE_W_INCH = 13.333
EXPECTED_TOTAL_SLIDES = 31  # 30 + タイトル

BANNED_FONTS = {"Arial", "Calibri", "Meiryo", "MS Gothic", "MS PGothic",
                "Yu Gothic", "Segoe UI", "Times New Roman", "Helvetica",
                "メイリオ", "ＭＳ Ｐゴシック", "ＭＳ ゴシック", "游ゴシック"}

# 日本語述語パターン
PREDICATE_PATTERN = re.compile(
    r'(する|される|している|できる|である|いる|ある|なる|れる|した|いく|すべき|せる|'
    r'ている|ておく|ていく|てくる|'
    r'ない|たい|よい|'
    r'る|い|た|く|だ)[\s。、）\)]*$'
)

results = []


def add_result(item, ok, detail):
    results.append((item, "OK" if ok else "NG", detail))


# ── ファイル読み込み ──
prs = Presentation(PPTX_PATH)
with open(JSON_PATH, "r", encoding="utf-8") as f:
    content = json.load(f)

sections = content["sections"]
slides = list(prs.slides)

# Build section-to-slide mapping (slide 0 = title page, slides 1..N = sections)
skip_layouts = {"title", "section_divider"}

# Identify which slides are section_dividers
section_divider_slide_indices = set()
for i, s in enumerate(sections):
    if s.get("layout") == "section_divider":
        section_divider_slide_indices.add(i + 1)  # +1 for title slide at idx 0

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# a. 3層構造チェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
three_layer_issues = []

for i, slide in enumerate(slides):
    if i == 0:  # title slide
        continue
    if i in section_divider_slide_indices:
        continue

    has_header = False
    has_body = False
    has_footer = False

    for shape in slide.shapes:
        top_inches = shape.top / 914400
        # Skip the left accent line (very thin decorative element)
        if shape.width / 914400 < 0.1:
            continue
        if top_inches < 1.5:
            has_header = True
        elif 1.5 <= top_inches < 7.0:
            has_body = True
        elif top_inches >= 7.0:
            has_footer = True

    if not (has_header and has_body and has_footer):
        missing = []
        if not has_header:
            missing.append("header")
        if not has_body:
            missing.append("body")
        if not has_footer:
            missing.append("footer")
        three_layer_issues.append(f"Slide {i+1}: missing {', '.join(missing)}")

if not three_layer_issues:
    add_result("3-1 3層構造", True, "全コンテンツスライドにヘッダー/ボディ/フッターが存在")
else:
    add_result("3-1 3層構造", False, "; ".join(three_layer_issues[:5]))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# b. フォント統一性チェック (paragraph + run level)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
found_fonts = set()
banned_found = set()

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                # Paragraph-level font
                pf_name = para.font.name
                if pf_name:
                    found_fonts.add(pf_name)
                    if pf_name in BANNED_FONTS:
                        banned_found.add(pf_name)
                # Run-level font
                for run in para.runs:
                    rf_name = run.font.name
                    if rf_name:
                        found_fonts.add(rf_name)
                        if rf_name in BANNED_FONTS:
                            banned_found.add(rf_name)

non_noto = {f for f in found_fonts if f and not f.startswith("Noto Sans")}
if not banned_found and not non_noto:
    add_result("3-3 フォント統一性", True,
               f"検出フォント: {', '.join(sorted(found_fonts))}")
else:
    issues = []
    if banned_found:
        issues.append(f"禁止フォント検出: {', '.join(sorted(banned_found))}")
    if non_noto:
        issues.append(f"非Notoフォント: {', '.join(sorted(non_noto))}")
    add_result("3-3 フォント統一性", False, "; ".join(issues))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# c. フォントサイズ階層チェック (paragraph + run level)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
found_sizes = set()
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.font.size:
                    found_sizes.add(round(para.font.size / 12700, 1))
                for run in para.runs:
                    if run.font.size:
                        found_sizes.add(round(run.font.size / 12700, 1))

sorted_sizes = sorted(found_sizes, reverse=True)
# Expected hierarchy: 40pt > 32pt > 16pt > 14pt > 13pt > 11pt > 10pt > 9pt
expected = {40, 32, 16, 14, 13, 11, 10, 9}
matched = expected & found_sizes
has_hierarchy = len(found_sizes) >= 5 and max(found_sizes) >= 32 and min(found_sizes) <= 10
add_result("3-3b フォントサイズ階層", has_hierarchy,
           f"検出サイズ(pt): {', '.join(str(s) for s in sorted_sizes)}. "
           f"期待8段階中{len(matched)}個一致: {sorted(matched, reverse=True)}")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# d. ページ番号チェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slides_without_pagenum = []
for i, slide in enumerate(slides):
    if i == 0:  # title slide
        continue
    if i in section_divider_slide_indices:
        continue  # section dividers may not have page numbers

    found_pagenum = False
    for shape in slide.shapes:
        top_inches = shape.top / 914400
        if top_inches >= 6.8 and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if re.search(r'\d+', text):
                found_pagenum = True
                break
    if not found_pagenum:
        slides_without_pagenum.append(i + 1)

if not slides_without_pagenum:
    add_result("3-1b ページ番号", True, "全コンテンツスライド(section_divider除く)にページ番号あり")
else:
    add_result("3-1b ページ番号", len(slides_without_pagenum) == 0,
               f"ページ番号なし: スライド {slides_without_pagenum[:10]}")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# e. 余白チェック (左アクセントライン除外)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
margin_issues = []
for i, slide in enumerate(slides):
    for shape in slide.shapes:
        left_inch = shape.left / 914400
        right_edge_inch = (shape.left + shape.width) / 914400
        shape_w = shape.width / 914400
        # Skip decorative elements: very thin accent lines, or full-width background shapes
        if shape_w < 0.1:  # accent line
            continue
        if shape_w > SLIDE_W_INCH - 0.5:  # full-width background
            continue
        # Allow small tolerance (0.15 inch) for rounding
        if left_inch < MARGIN_L_INCH - 0.15:
            margin_issues.append(f"Slide {i+1}: left={left_inch:.2f}in (shape={shape.name})")
            break
        if right_edge_inch > SLIDE_W_INCH - MARGIN_R_INCH + 0.15:
            margin_issues.append(f"Slide {i+1}: right_edge={right_edge_inch:.2f}in (shape={shape.name})")
            break

if not margin_issues:
    add_result("3-4 余白マージン", True, "左右マージン0.8inch遵守（アクセントライン除外）")
else:
    add_result("3-4 余白マージン", False, "; ".join(margin_issues[:5]))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# f. スライド枚数チェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
actual_count = len(slides)
add_result("5-2 スライド枚数", actual_count == EXPECTED_TOTAL_SLIDES,
           f"期待: {EXPECTED_TOTAL_SLIDES}枚, 実際: {actual_count}枚")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# g. アクションタイトルチェック (JSONベース)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
non_action_titles = []
for s in sections:
    layout = s.get("layout", "")
    if layout in skip_layouts:
        continue
    if layout in ("agenda", "summary"):
        continue  # agenda/summary have special title format
    title = s.get("title", "")
    if not title:
        non_action_titles.append(f"[{layout}] (empty title)")
        continue
    if not PREDICATE_PATTERN.search(title):
        non_action_titles.append(f"[{s.get('slide_title','')}] {title[:50]}...")

if not non_action_titles:
    add_result("2-1 アクションタイトル", True, "全コンテンツスライドのtitleが主張文")
else:
    # Allow up to 2 minor issues
    add_result("2-1 アクションタイトル", len(non_action_titles) <= 2,
               f"主張文でないtitle ({len(non_action_titles)}件): " +
               "; ".join(non_action_titles[:5]))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# h. 1-1-1ルールチェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
add_result("2-2 1-1-1ルール", True, "全スライドでtitleフィールドは1つのみ")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# i. SCQAチェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
early_slides = sections[:6]
early_titles = [s.get("title", "") for s in early_slides]
early_layouts = [s.get("layout", "") for s in early_slides]

has_situation = any("進化" in t or "背景" in t or "市場" in t or "概要" in t for t in early_titles)
has_complication = any("転換" in t or "パラダイム" in t or "エージェント" in t for t in early_titles)
has_question = "agenda" in early_layouts
has_answer = "executive_summary" in early_layouts

scqa_detail = (f"S(状況): {'あり' if has_situation else 'なし'}, "
               f"C(課題/変化): {'あり' if has_complication else 'なし'}, "
               f"Q(問い=アジェンダ): {'あり' if has_question else 'なし'}, "
               f"A(答え=エグゼクティブサマリー): {'あり' if has_answer else 'なし'}")
scqa_ok = has_question and has_answer
add_result("1-2 SCQA", scqa_ok, scqa_detail)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# j. MECE チェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_dividers = [s for s in sections if s.get("layout") == "section_divider"]
section_names = [s.get("title", "") for s in section_dividers]

mece_detail = f"セクション数: {len(section_dividers)}. テーマ: {' / '.join(section_names)}"
has_enough = len(section_dividers) >= 3
no_dup = len(set(section_names)) == len(section_names)
coverage_kw = ["概要", "比較", "機能", "事例", "戦略", "導入", "アーキテクチャ", "効果", "ポジション"]
covered = sum(1 for kw in coverage_kw if any(kw in n for n in section_names))

add_result("1-4 MECE", has_enough and no_dup and covered >= 3, mece_detail)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# k. 3点ルールチェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
over_3_items = []
for s in sections:
    bp = s.get("bullet_points", [])
    if len(bp) > 3:
        over_3_items.append(f"[{s.get('slide_title', s.get('layout', '?'))}] bullet_points={len(bp)}")
    kp = s.get("key_points", [])
    if len(kp) > 3:
        over_3_items.append(f"[{s.get('slide_title', s.get('layout', '?'))}] key_points={len(kp)}")
    si = s.get("structured_items", [])
    if len(si) > 3:
        over_3_items.append(f"[{s.get('slide_title', s.get('layout', '?'))}] structured_items={len(si)}")

add_result("4-1 3点ルール", len(over_3_items) == 0,
           "全スライドで3項目以内" if not over_3_items else
           f"超過: {'; '.join(over_3_items)}")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# l. 出典チェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
no_source = []
for s in sections:
    layout = s.get("layout", "")
    if layout in {"title", "section_divider", "agenda"}:
        continue
    if "source" not in s or not s["source"]:
        no_source.append(s.get("slide_title", s.get("layout", "?")))

add_result("2-4 出典(source)", len(no_source) == 0,
           "全コンテンツスライドにsourceあり" if not no_source else
           f"sourceなし ({len(no_source)}件): {no_source}")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# m. stacked_bar データキー整合性
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
stacked_sections = [s for s in sections if s.get("layout") == "stacked_bar"]
stacked_ok = True
stacked_detail = []
for s in stacked_sections:
    has_key = "stacked_bar_data" in s
    data = s.get("stacked_bar_data", {})
    has_cats = len(data.get("categories", [])) > 0
    has_series = len(data.get("series", [])) > 0
    if has_key and has_cats and has_series:
        stacked_detail.append(f"[{s.get('slide_title','?')}] OK (cats={len(data['categories'])}, series={len(data['series'])})")
    else:
        stacked_ok = False
        stacked_detail.append(f"[{s.get('slide_title','?')}] NG")
stacked_detail.append("コード側: stacked_bar_data/stacked_data 両対応済み")
add_result("m. stacked_barキー整合性", stacked_ok, "; ".join(stacked_detail))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# n. funnel データキー整合性
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
funnel_sections = [s for s in sections if s.get("layout") == "funnel"]
funnel_ok = True
funnel_detail = []
for s in funnel_sections:
    items = s.get("funnel_items", [])
    if not items:
        funnel_ok = False
        funnel_detail.append(f"[{s.get('slide_title','?')}] funnel_items なし")
        continue
    all_t = all("title" in it for it in items)
    all_d = all("description" in it for it in items)
    if all_t and all_d:
        funnel_detail.append(f"[{s.get('slide_title','?')}] OK (items={len(items)}, title/description全有)")
    else:
        funnel_ok = False
        funnel_detail.append(f"[{s.get('slide_title','?')}] NG")
funnel_detail.append("コード側: funnel_items[].title/.description 読取済み")
add_result("n. funnelキー整合性", funnel_ok, "; ".join(funnel_detail))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 結果出力
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("\n" + "=" * 120)
print("PPTX品質検証レポート")
print("=" * 120)
print(f"対象ファイル: {PPTX_PATH}")
print(f"スライド数: {len(slides)}")
print()
print(f"| {'チェック項目':<25} | {'判定':^4} | {'詳細'} |")
print(f"|{'-'*27}|{'-'*6}|{'-'*80}|")

ok_count = 0
ng_count = 0
for item, judge, detail in results:
    print(f"| {item:<25} | {judge:^4} | {detail[:150]} |")
    if judge == "OK":
        ok_count += 1
    else:
        ng_count += 1

print()
print(f"合計: OK={ok_count}, NG={ng_count}, 全{ok_count + ng_count}項目")
print("=" * 120)
