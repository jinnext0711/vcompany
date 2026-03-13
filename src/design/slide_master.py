"""スライドマスタ（テンプレート）生成モジュール

python-pptx を使い、DesignConfig に準拠したスライドレイアウトを
動的に生成する。各レイアウトは「ヘッダー／ボディ／フッター」の
3 層構造を基本とする。
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from src.design.builder import DesignConfig


# ---------------------------------------------------------------------------
# ヘルパー
# ---------------------------------------------------------------------------

def _rgb(hex_color: str) -> RGBColor:
    """6 桁 hex 文字列を RGBColor に変換する。"""
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _apply_font(run, font_cfg: dict, color_hex: str | None = None):
    """Run オブジェクトにフォント設定を適用する。"""
    run.font.name = font_cfg["name"]
    run.font.size = Pt(font_cfg["size"])
    run.font.bold = font_cfg.get("bold", False)
    if color_hex:
        run.font.color.rgb = _rgb(color_hex)


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_cfg: dict,
    color_hex: str,
    alignment: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    bg_hex: str | None = None,
):
    """テキストボックスを追加して書式を設定する。"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # 垂直位置
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _apply_font(run, font_cfg, color_hex)

    return txBox


def _add_rect_bg(slide, left, top, width, height, fill_hex):
    """背景用の塗りつぶし矩形を追加する。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()  # 枠線なし
    # 最背面に移動（xml 操作）
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def _add_footer_elements(slide, cfg: DesignConfig, layout: dict):
    """フッター（出典プレースホルダー + ページ番号）を追加する。"""
    sw = layout["slide_width"]
    sh = layout["slide_height"]
    ml = layout["margin_left"]
    mr = layout["margin_right"]
    footer_h = 0.4
    footer_top = sh - footer_h

    body_w = sw - ml - mr

    # 区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(ml), Inches(footer_top - 0.05), Inches(body_w), Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(cfg.colors["border"])
    line.line.fill.background()

    # 出典テキスト
    _add_textbox(
        slide,
        left=ml,
        top=footer_top,
        width=body_w * 0.75,
        height=footer_h,
        text="出典: ",
        font_cfg=cfg.get_font("caption"),
        color_hex=cfg.colors["light_text"],
        alignment=PP_ALIGN.LEFT,
    )

    # ページ番号プレースホルダー
    _add_textbox(
        slide,
        left=ml + body_w * 0.75,
        top=footer_top,
        width=body_w * 0.25,
        height=footer_h,
        text="##",
        font_cfg=cfg.get_font("caption"),
        color_hex=cfg.colors["light_text"],
        alignment=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# SlideMasterBuilder
# ---------------------------------------------------------------------------

class SlideMasterBuilder:
    """DesignConfig に基づくスライドレイアウト群を生成する。

    python-pptx ではスライドマスタの XML を直接操作するのが困難なため、
    ブランクレイアウト上にプレースホルダー相当のシェイプ群を配置する
    「テンプレートスライド」方式で実装する。

    Usage::

        builder = SlideMasterBuilder()
        prs = builder.build_presentation()
        prs.save("template.pptx")
    """

    # ヘッダー既定値
    HEADER_HEIGHT = 1.3   # inches
    FOOTER_HEIGHT = 0.4   # inches

    def __init__(self, config: DesignConfig | None = None):
        self.cfg = config or DesignConfig()
        self.layout = DesignConfig.LAYOUT
        self._prs: Presentation | None = None

    # ------------------------------------------------------------------
    # 内部ユーティリティ
    # ------------------------------------------------------------------

    def _new_prs(self) -> Presentation:
        """16:9 の空 Presentation を作成する。"""
        prs = Presentation()
        prs.slide_width = Inches(self.layout["slide_width"])
        prs.slide_height = Inches(self.layout["slide_height"])
        return prs

    def _blank_layout(self, prs: Presentation):
        """ブランクスライドレイアウトを返す。"""
        return prs.slide_layouts[6]  # 通常 index 6 がブランク

    def _body_top(self) -> float:
        return self.layout["margin_top"] + self.HEADER_HEIGHT + self.layout["content_gap"]

    def _body_height(self) -> float:
        return (
            self.layout["slide_height"]
            - self._body_top()
            - self.FOOTER_HEIGHT
            - self.layout["content_gap"]
        )

    def _body_width(self) -> float:
        return self.layout["slide_width"] - self.layout["margin_left"] - self.layout["margin_right"]

    def _add_header(self, slide, title_text: str = "アクションタイトル"):
        """濃紺背景 + 白文字 20pt のヘッダーバーを追加する。"""
        sw = self.layout["slide_width"]
        _add_rect_bg(slide, 0, 0, sw, self.HEADER_HEIGHT, self.cfg.colors["primary"])

        header_font = {
            "name": self.cfg.get_font("subheading")["name"],
            "size": 20,
            "bold": True,
        }
        _add_textbox(
            slide,
            left=self.layout["margin_left"],
            top=0.25,
            width=sw - self.layout["margin_left"] - self.layout["margin_right"],
            height=self.HEADER_HEIGHT - 0.25,
            text=title_text,
            font_cfg=header_font,
            color_hex="FFFFFF",
            alignment=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    def _add_footer(self, slide):
        _add_footer_elements(slide, self.cfg, self.layout)

    # ------------------------------------------------------------------
    # レイアウト生成メソッド
    # ------------------------------------------------------------------

    def create_title_layout(self, prs: Presentation) -> object:
        """タイトルスライド — 中央揃えのタイトル + サブタイトル。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        sw = self.layout["slide_width"]
        sh = self.layout["slide_height"]

        # 全面濃紺背景
        _add_rect_bg(slide, 0, 0, sw, sh, self.cfg.colors["primary"])

        # タイトル
        _add_textbox(
            slide,
            left=1.5,
            top=sh * 0.25,
            width=sw - 3.0,
            height=1.5,
            text="プレゼンテーションタイトル",
            font_cfg=self.cfg.get_font("title"),
            color_hex="FFFFFF",
            alignment=PP_ALIGN.CENTER,
        )

        # サブタイトル
        _add_textbox(
            slide,
            left=2.0,
            top=sh * 0.25 + 1.8,
            width=sw - 4.0,
            height=0.8,
            text="サブタイトル / 日付 / 発表者名",
            font_cfg=self.cfg.get_font("body"),
            color_hex=self.cfg.colors["border"],
            alignment=PP_ALIGN.CENTER,
        )

        # 区切りライン（アクセント）
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(sw / 2 - 1.5), Inches(sh * 0.25 + 1.55),
            Inches(3.0), Inches(0.04),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(self.cfg.colors["highlight"])
        line.line.fill.background()

        return slide

    def create_section_divider_layout(self, prs: Presentation) -> object:
        """セクション区切りスライド — セクション番号 + セクション名。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        sw = self.layout["slide_width"]
        sh = self.layout["slide_height"]

        # 背景: 濃紺
        _add_rect_bg(slide, 0, 0, sw, sh, self.cfg.colors["primary"])

        # セクション番号
        num_font = {"name": self.cfg.get_font("data")["name"], "size": 72, "bold": True}
        _add_textbox(
            slide,
            left=self.layout["margin_left"],
            top=sh * 0.25,
            width=2.0,
            height=1.5,
            text="01",
            font_cfg=num_font,
            color_hex=self.cfg.colors["highlight"],
            alignment=PP_ALIGN.LEFT,
        )

        # セクション名
        _add_textbox(
            slide,
            left=self.layout["margin_left"] + 2.2,
            top=sh * 0.30,
            width=sw - self.layout["margin_left"] - self.layout["margin_right"] - 2.2,
            height=1.2,
            text="セクションタイトル",
            font_cfg=self.cfg.get_font("heading"),
            color_hex="FFFFFF",
            alignment=PP_ALIGN.LEFT,
        )

        return slide

    def create_content_layout(self, prs: Presentation) -> object:
        """コンテンツスライド — ヘッダー + ボディ + フッター 3 層構造。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        ml = self.layout["margin_left"]

        self._add_header(slide)

        # ボディ: テキスト領域
        _add_textbox(
            slide,
            left=ml,
            top=self._body_top(),
            width=self._body_width(),
            height=self._body_height(),
            text="本文テキスト\n— ポイント1\n— ポイント2\n— ポイント3",
            font_cfg=self.cfg.get_font("body"),
            color_hex=self.cfg.colors["text"],
            alignment=PP_ALIGN.LEFT,
        )

        self._add_footer(slide)
        return slide

    def create_two_column_layout(self, prs: Presentation) -> object:
        """2 カラムスライド — 左右均等分割。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        ml = self.layout["margin_left"]
        gap = self.layout["content_gap"]
        col_w = (self._body_width() - gap) / 2

        self._add_header(slide)

        # 左カラム
        _add_textbox(
            slide,
            left=ml,
            top=self._body_top(),
            width=col_w,
            height=self._body_height(),
            text="左カラム内容",
            font_cfg=self.cfg.get_font("body"),
            color_hex=self.cfg.colors["text"],
            alignment=PP_ALIGN.LEFT,
        )

        # 右カラム
        _add_textbox(
            slide,
            left=ml + col_w + gap,
            top=self._body_top(),
            width=col_w,
            height=self._body_height(),
            text="右カラム内容",
            font_cfg=self.cfg.get_font("body"),
            color_hex=self.cfg.colors["text"],
            alignment=PP_ALIGN.LEFT,
        )

        self._add_footer(slide)
        return slide

    def create_chart_layout(self, prs: Presentation) -> object:
        """図表中心スライド — ヘッダー + 図表領域 + 注釈 + フッター。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        ml = self.layout["margin_left"]
        body_w = self._body_width()
        body_h = self._body_height()

        self._add_header(slide)

        # 図表プレースホルダー（ライトグレー背景の矩形）
        chart_h = body_h * 0.78
        chart_shape = _add_rect_bg(
            slide, ml, self._body_top(), body_w, chart_h,
            self.cfg.colors["light_bg"],
        )

        # 中央に「図表エリア」ラベル
        _add_textbox(
            slide,
            left=ml + body_w * 0.3,
            top=self._body_top() + chart_h * 0.4,
            width=body_w * 0.4,
            height=0.6,
            text="図表エリア",
            font_cfg=self.cfg.get_font("subheading"),
            color_hex=self.cfg.colors["light_text"],
            alignment=PP_ALIGN.CENTER,
        )

        # 注釈行
        _add_textbox(
            slide,
            left=ml,
            top=self._body_top() + chart_h + 0.1,
            width=body_w,
            height=body_h - chart_h - 0.1,
            text="注: ",
            font_cfg=self.cfg.get_font("caption"),
            color_hex=self.cfg.colors["light_text"],
            alignment=PP_ALIGN.LEFT,
        )

        self._add_footer(slide)
        return slide

    def create_kpi_layout(self, prs: Presentation) -> object:
        """数値ハイライト (KPI) スライド — 大数字 + ラベルを 3 列配置。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        ml = self.layout["margin_left"]
        gap = self.layout["content_gap"]
        body_w = self._body_width()
        body_h = self._body_height()
        col_count = 3
        col_w = (body_w - gap * (col_count - 1)) / col_count

        self._add_header(slide, title_text="主要指標")

        for i in range(col_count):
            col_left = ml + i * (col_w + gap)
            kpi_top = self._body_top() + body_h * 0.15

            # 数値
            _add_textbox(
                slide,
                left=col_left,
                top=kpi_top,
                width=col_w,
                height=1.2,
                text=f"0{i + 1}",
                font_cfg=self.cfg.get_font("data"),
                color_hex=self.cfg.colors["highlight"],
                alignment=PP_ALIGN.CENTER,
            )

            # ラベル
            _add_textbox(
                slide,
                left=col_left,
                top=kpi_top + 1.3,
                width=col_w,
                height=0.6,
                text=f"指標{i + 1}の説明",
                font_cfg=self.cfg.get_font("caption"),
                color_hex=self.cfg.colors["text"],
                alignment=PP_ALIGN.CENTER,
            )

        self._add_footer(slide)
        return slide

    def create_summary_layout(self, prs: Presentation) -> object:
        """まとめスライド — 主張 + 箇条書き要約。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        ml = self.layout["margin_left"]
        body_w = self._body_width()
        body_h = self._body_height()

        self._add_header(slide, title_text="まとめ")

        # 主張（太字・大きめ）
        _add_textbox(
            slide,
            left=ml,
            top=self._body_top(),
            width=body_w,
            height=1.0,
            text="結論 / 主張を1文で記載",
            font_cfg=self.cfg.get_font("heading"),
            color_hex=self.cfg.colors["primary"],
            alignment=PP_ALIGN.LEFT,
        )

        # 箇条書き要約
        _add_textbox(
            slide,
            left=ml,
            top=self._body_top() + 1.2,
            width=body_w,
            height=body_h - 1.2,
            text="— 要点1\n— 要点2\n— 要点3",
            font_cfg=self.cfg.get_font("body"),
            color_hex=self.cfg.colors["text"],
            alignment=PP_ALIGN.LEFT,
        )

        self._add_footer(slide)
        return slide

    def create_agenda_layout(self, prs: Presentation) -> object:
        """アジェンダスライド — 番号付きリスト。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        ml = self.layout["margin_left"]
        body_w = self._body_width()
        body_h = self._body_height()

        self._add_header(slide, title_text="アジェンダ")

        item_count = 5
        item_h = min(body_h / item_count, 0.9)

        for i in range(item_count):
            item_top = self._body_top() + i * item_h

            # 番号
            num_font = {"name": self.cfg.get_font("heading")["name"], "size": 24, "bold": True}
            _add_textbox(
                slide,
                left=ml,
                top=item_top,
                width=0.6,
                height=item_h,
                text=f"{i + 1:02d}",
                font_cfg=num_font,
                color_hex=self.cfg.colors["highlight"],
                alignment=PP_ALIGN.RIGHT,
            )

            # テキスト
            _add_textbox(
                slide,
                left=ml + 0.8,
                top=item_top,
                width=body_w - 0.8,
                height=item_h,
                text=f"議題{i + 1}",
                font_cfg=self.cfg.get_font("body"),
                color_hex=self.cfg.colors["text"],
                alignment=PP_ALIGN.LEFT,
            )

            # 区切り線（最後以外）
            if i < item_count - 1:
                sep = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(ml + 0.8), Inches(item_top + item_h - 0.04),
                    Inches(body_w - 0.8), Inches(0.015),
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = _rgb(self.cfg.colors["border"])
                sep.line.fill.background()

        self._add_footer(slide)
        return slide

    def create_executive_summary_layout(self, prs: Presentation) -> object:
        """エグゼクティブサマリー — 結論 + 根拠 (2 カラム) + ネクストステップ。"""
        slide = prs.slides.add_slide(self._blank_layout(prs))
        ml = self.layout["margin_left"]
        gap = self.layout["content_gap"]
        body_w = self._body_width()
        body_h = self._body_height()

        self._add_header(slide, title_text="エグゼクティブサマリー")

        # 結論帯（ハイライト背景）
        conclusion_h = 0.9
        _add_rect_bg(
            slide, ml - 0.1, self._body_top(),
            body_w + 0.2, conclusion_h,
            self.cfg.colors["light_bg"],
        )
        _add_textbox(
            slide,
            left=ml,
            top=self._body_top() + 0.1,
            width=body_w,
            height=conclusion_h - 0.2,
            text="結論: 主要な結論を1-2文で記載",
            font_cfg=self.cfg.get_font("subheading"),
            color_hex=self.cfg.colors["primary"],
            alignment=PP_ALIGN.LEFT,
        )

        # 根拠（2 カラム）
        evidence_top = self._body_top() + conclusion_h + gap
        evidence_h = body_h - conclusion_h - gap - 1.2
        col_w = (body_w - gap) / 2

        _add_textbox(
            slide,
            left=ml,
            top=evidence_top,
            width=col_w,
            height=evidence_h,
            text="根拠 1\n— データポイントA\n— データポイントB",
            font_cfg=self.cfg.get_font("body"),
            color_hex=self.cfg.colors["text"],
            alignment=PP_ALIGN.LEFT,
        )

        _add_textbox(
            slide,
            left=ml + col_w + gap,
            top=evidence_top,
            width=col_w,
            height=evidence_h,
            text="根拠 2\n— データポイントC\n— データポイントD",
            font_cfg=self.cfg.get_font("body"),
            color_hex=self.cfg.colors["text"],
            alignment=PP_ALIGN.LEFT,
        )

        # ネクストステップ
        next_top = evidence_top + evidence_h + 0.15
        _add_textbox(
            slide,
            left=ml,
            top=next_top,
            width=body_w,
            height=0.8,
            text="ネクストステップ: 推奨アクションを記載",
            font_cfg=self.cfg.get_font("subheading"),
            color_hex=self.cfg.colors["highlight"],
            alignment=PP_ALIGN.LEFT,
        )

        self._add_footer(slide)
        return slide

    # ------------------------------------------------------------------
    # ビルド
    # ------------------------------------------------------------------

    def build_presentation(self) -> Presentation:
        """全レイアウトを含むベース Presentation オブジェクトを返す。

        スライド順:
          0. タイトル
          1. アジェンダ
          2. セクション区切り
          3. コンテンツ
          4. 2 カラム
          5. 図表中心
          6. KPI
          7. エグゼクティブサマリー
          8. まとめ
        """
        prs = self._new_prs()

        self.create_title_layout(prs)
        self.create_agenda_layout(prs)
        self.create_section_divider_layout(prs)
        self.create_content_layout(prs)
        self.create_two_column_layout(prs)
        self.create_chart_layout(prs)
        self.create_kpi_layout(prs)
        self.create_executive_summary_layout(prs)
        self.create_summary_layout(prs)

        self._prs = prs
        return prs
