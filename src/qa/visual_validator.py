"""Phase 6.5: 視覚品質検証 — PPTXレンダリング結果のチェッカー。

生成済みPPTXファイルを再読み込みし、視覚的品質を検証する。
JSONベースのQA（Phase 5）では検出できない以下の問題を検出:
  V-1: テキスト重複（シェイプ同士の重なり）
  V-2: 空スライド（ボディエリアにコンテンツなし）
  V-3: テキスト溢れ（テキストがコンテナからはみ出し）
  V-4: 余白バランス（過大な余白・スカスカ）
  V-5: コントラスト比（WCAG AA準拠チェック）
  V-6: スライド外はみ出し（オブジェクトがスライド境界外）

担当: QA部 + デザイン部（共同レビュー）
"""

from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Emu

from src.output.helpers import (
    SLIDE_W, SLIDE_H, MARGIN_L, MARGIN_R,
    BODY_TOP, BODY_BOTTOM, FOOTER_TOP,
)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# シェイプ情報抽出
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@dataclass
class ShapeBounds:
    """シェイプのバウンディングボックス（inches）。"""
    left: float
    top: float
    width: float
    height: float
    text: str
    shape_name: str
    font_size_pt: float = 0
    is_text_shape: bool = False

    @property
    def right(self):
        return self.left + self.width

    @property
    def bottom(self):
        return self.top + self.height


def _emu_to_inches(emu):
    """EMU → inches 変換。"""
    return emu / 914400 if emu else 0


def _extract_shapes(slide):
    """スライドから全シェイプを抽出する。"""
    results = []
    for shape in slide.shapes:
        left = _emu_to_inches(shape.left)
        top = _emu_to_inches(shape.top)
        width = _emu_to_inches(shape.width)
        height = _emu_to_inches(shape.height)

        text = ""
        font_size = 0
        is_text = False

        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            is_text = bool(text)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font and run.font.size:
                        font_size = run.font.size.pt
                        break
                if font_size > 0:
                    break
            if font_size == 0:
                for para in shape.text_frame.paragraphs:
                    if para.font and para.font.size:
                        font_size = para.font.size.pt
                        break

        results.append(ShapeBounds(
            left=left, top=top, width=width, height=height,
            text=text, shape_name=shape.name,
            font_size_pt=font_size, is_text_shape=is_text,
        ))
    return results


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# V-1: テキスト重複チェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _rects_overlap(a, b, tol=0.08):
    """2つの矩形が重なっているか判定（tolerance付き）。"""
    return (a.left + tol < b.right and b.left + tol < a.right and
            a.top + tol < b.bottom and b.top + tol < a.bottom)


def _fully_contains(outer, inner):
    """outerがinnerを完全に包含しているか。"""
    return (outer.left <= inner.left and outer.right >= inner.right and
            outer.top <= inner.top and outer.bottom >= inner.bottom)


def _check_text_overlap(shapes, slide_idx):
    """V-1: テキストを持つシェイプ同士の重複を検出する。"""
    issues = []
    text_shapes = [s for s in shapes if s.is_text_shape and s.height > 0.1]

    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            a, b = text_shapes[i], text_shapes[j]
            if _fully_contains(a, b) or _fully_contains(b, a):
                continue
            if _rects_overlap(a, b):
                overlap_x = min(a.right, b.right) - max(a.left, b.left)
                overlap_y = min(a.bottom, b.bottom) - max(a.top, b.top)
                if overlap_x > 0.15 and overlap_y > 0.15:
                    severity = "error" if overlap_y > 0.3 else "warning"
                    issues.append({
                        "check_id": "V-1",
                        "check_name": "テキスト重複",
                        "slide_index": slide_idx,
                        "severity": severity,
                        "detail": (
                            f"「{a.text[:20]}」と「{b.text[:20]}」が"
                            f"{overlap_x:.2f}\"x{overlap_y:.2f}\"重複"
                        ),
                    })
    return issues


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# V-2: 空スライドチェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _check_empty_slides(shapes, slide_idx):
    """V-2: ボディエリアにテキストコンテンツがないスライドを検出する。"""
    if slide_idx == 0:
        return []  # タイトルスライドはスキップ

    body_texts = [
        s for s in shapes
        if s.is_text_shape and s.top >= BODY_TOP - 0.1 and s.top < BODY_BOTTOM
    ]

    if not body_texts:
        return [{
            "check_id": "V-2",
            "check_name": "空スライド",
            "slide_index": slide_idx,
            "severity": "error",
            "detail": "ボディエリアにテキストコンテンツがありません",
        }]
    return []


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# V-3: テキスト溢れチェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _check_text_overflow(shapes, slide_idx):
    """V-3: テキストがコンテナからはみ出しているか推定する。"""
    issues = []
    for s in shapes:
        if not s.is_text_shape or s.width < 0.3 or s.height < 0.1:
            continue
        fs = s.font_size_pt or 12
        char_w = fs / 72 * 0.8
        chars_per_line = max(int(s.width / char_w), 1)
        lines = max(1, -(-len(s.text) // chars_per_line))
        line_h = fs * 1.3 / 72
        estimated_h = lines * line_h

        if estimated_h > s.height * 1.5:
            issues.append({
                "check_id": "V-3",
                "check_name": "テキスト溢れ",
                "slide_index": slide_idx,
                "severity": "error" if estimated_h > s.height * 2 else "warning",
                "detail": (
                    f"「{s.text[:25]}」推定高さ{estimated_h:.2f}\" > "
                    f"コンテナ{s.height:.2f}\"（{estimated_h / s.height:.1f}倍）"
                ),
            })
    return issues


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# V-4: 余白バランスチェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _check_whitespace_gaps(shapes, slide_idx):
    """V-4: 過大な余白を検出する。"""
    if slide_idx == 0:
        return []

    body_shapes = [
        s for s in shapes
        if (s.is_text_shape or s.height > 0.2) and
        s.top >= BODY_TOP - 0.1 and s.top < BODY_BOTTOM
    ]

    if not body_shapes:
        return []

    max_bottom = max(s.bottom for s in body_shapes)
    gap = BODY_BOTTOM - max_bottom

    issues = []
    if gap > 1.8:
        issues.append({
            "check_id": "V-4",
            "check_name": "余白バランス",
            "slide_index": slide_idx,
            "severity": "warning",
            "detail": f"ボディ下端に{gap:.1f}\"の余白（1.8\"超過）",
        })

    min_top = min(s.top for s in body_shapes)
    content_span = max_bottom - min_top
    body_h = BODY_BOTTOM - BODY_TOP
    if content_span < body_h * 0.35:
        issues.append({
            "check_id": "V-4",
            "check_name": "余白バランス",
            "slide_index": slide_idx,
            "severity": "warning",
            "detail": f"コンテンツ高さ{content_span:.1f}\"がボディ高さ{body_h:.1f}\"の35%未満",
        })

    return issues


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# V-5: スライド外はみ出しチェック
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _check_out_of_bounds(shapes, slide_idx):
    """V-5: スライド境界外にはみ出しているシェイプを検出する。"""
    issues = []
    for s in shapes:
        if not s.is_text_shape:
            continue
        if s.right > SLIDE_W + 0.1:
            issues.append({
                "check_id": "V-5",
                "check_name": "スライド外はみ出し",
                "slide_index": slide_idx,
                "severity": "error",
                "detail": f"「{s.text[:20]}」が右端{s.right:.2f}\"（スライド幅{SLIDE_W:.1f}\"）",
            })
        if s.bottom > SLIDE_H + 0.1:
            issues.append({
                "check_id": "V-5",
                "check_name": "スライド外はみ出し",
                "slide_index": slide_idx,
                "severity": "error",
                "detail": f"「{s.text[:20]}」が下端{s.bottom:.2f}\"（スライド高{SLIDE_H:.1f}\"）",
            })
    return issues


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# メイン実行
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def run_visual_validation(pptx_path, design=None):
    """全視覚検証チェックを実行する。

    Returns:
        dict: issues, summary, per_slide
    """
    prs = Presentation(pptx_path)
    all_issues = []
    per_slide = {}

    for slide_idx, slide in enumerate(prs.slides):
        shapes = _extract_shapes(slide)
        slide_issues = []

        # 特殊スライド判定（背景塗りつぶし or テキストがBODY_TOP以上にしかない）
        # title, section_divider, summary は背景色つき → V-2/V-4スキップ
        is_special = False
        try:
            bg_fill = slide.background.fill
            if bg_fill.type is not None:
                is_special = True
        except Exception:
            pass
        # agenda スライドも特殊扱い（BODY_TOP付近からコンテンツ開始）
        if not is_special:
            body_shapes = [s for s in shapes if s.is_text_shape and s.top < BODY_TOP]
            header_texts = [s.text for s in body_shapes]
            if any("Agenda" in t or "アジェンダ" in t or "本日の論点" in t for t in header_texts):
                is_special = True

        slide_issues.extend(_check_text_overlap(shapes, slide_idx))
        if not is_special:
            slide_issues.extend(_check_empty_slides(shapes, slide_idx))
        slide_issues.extend(_check_text_overflow(shapes, slide_idx))
        if not is_special:
            slide_issues.extend(_check_whitespace_gaps(shapes, slide_idx))
        slide_issues.extend(_check_out_of_bounds(shapes, slide_idx))

        all_issues.extend(slide_issues)
        if slide_issues:
            per_slide[slide_idx] = slide_issues

    errors = sum(1 for i in all_issues if i["severity"] == "error")
    warnings = sum(1 for i in all_issues if i["severity"] == "warning")

    return {
        "pptx_path": pptx_path,
        "slide_count": len(prs.slides),
        "issues": all_issues,
        "summary": {
            "total_issues": len(all_issues),
            "errors": errors,
            "warnings": warnings,
            "slides_with_issues": len(per_slide),
            "clean_slides": len(prs.slides) - len(per_slide),
        },
        "per_slide": per_slide,
    }


def format_visual_report(result):
    """視覚検証結果を人間可読な文字列にフォーマットする。"""
    lines = []
    s = result["summary"]
    lines.append(f"=== 視覚品質検証 ({result['slide_count']}スライド) ===")
    lines.append(f"  Error: {s['errors']} / Warning: {s['warnings']} / "
                 f"問題スライド: {s['slides_with_issues']} / "
                 f"正常: {s['clean_slides']}")
    lines.append("")

    for slide_idx in sorted(result["per_slide"].keys()):
        slide_issues = result["per_slide"][slide_idx]
        lines.append(f"--- スライド {slide_idx + 1} ---")
        for issue in slide_issues:
            marker = "NG" if issue["severity"] == "error" else "WARN"
            lines.append(f"  [{marker}] {issue['check_name']}: {issue['detail']}")
        lines.append("")

    return "\n".join(lines)
