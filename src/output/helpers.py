"""共通描画ヘルパー関数群。

全レイアウト・図解レンダラーが共有する基本描画プリミティブ。
座標系は consulting-slide-master-spec.md（MBB/Big4統合仕様）に完全準拠。
"""

import math
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from src.design.builder import DesignConfig

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# レイアウト定数 — consulting-slide-master-spec 完全準拠
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SLIDE_W = 13.333          # 16:9 スライド幅 (inches)
SLIDE_H = 7.5             # 16:9 スライド高 (inches)
MARGIN_L = 0.50           # 左マージン (spec §3-3)
MARGIN_R = 0.50           # 右マージン (spec §3-3)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R  # = 12.333"

# ヘッダー3層構造 (spec §3-2, §5-1)
HEADER_TOP = 0.30          # アクションタイトル開始Y
HEADER_HEIGHT = 0.83       # アクションタイトル高さ（20pt×2行分）
ACCENT_LINE_TOP = 1.15     # アクセント区切り線Y
ACCENT_LINE_H = 0.017      # 区切り線の太さ（≒1.2pt）

# 互換エイリアス
HEADER_TITLE_TOP = 0.30    # スライドタイトル開始Y
HEADER_MSG_TOP = 0.62      # メッセージ開始Y（タイトルと重なるエリア内）
HEADER_BOTTOM = ACCENT_LINE_TOP  # ヘッダー下端 = 区切り線位置

# ボディエリア (spec §3-3)
BODY_TOP = 1.25            # ボディ開始Y
BODY_BOTTOM = 6.80         # ボディ下端Y
BODY_H = BODY_BOTTOM - BODY_TOP  # = 5.55"

# フッター (spec §3-3, §5-3)
FOOTER_TOP = 7.05          # フッター開始Y
FOOTER_H = 0.30            # フッター高さ

# グリッドシステム (spec §5-4)
GAP = 0.333                # カラム間ギャップ
COL2_W = (CONTENT_W - GAP) / 2                # = 6.0"
COL3_W = (CONTENT_W - GAP * 2) / 3            # = 3.889"
COL_1_3_W = (CONTENT_W - GAP) * 1 / 3         # = 3.889"  (1/3)
COL_2_3_W = (CONTENT_W - GAP) * 2 / 3         # = 8.111"  (2/3)

# 8pxグリッド (0.11 inches ≒ 8px at 72dpi)
G = 0.11

# フォントサイズ階層 (pt) — spec §7-3 FONT_SIZES 準拠
FS_TITLE = 36              # 表紙タイトル
FS_ACTION_TITLE = 20       # アクションタイトル（全コンテンツスライド）
FS_SECTION_TITLE = 32      # セクション区切りタイトル
FS_SECTION_NUMBER = 48     # セクション番号
FS_SUBTITLE = 18           # サブタイトル
FS_HEADING = 16            # セクション見出し（ボディ内）
FS_BODY = 14               # 本文
FS_BODY_SMALL = 12         # 小さい本文
FS_TABLE_HEADER = 11       # テーブルヘッダー
FS_TABLE_BODY = 11         # テーブル本文（最低可読サイズ）
FS_DATA_LABEL = 11         # データラベル・グラフラベル
FS_CAPTION = 11            # キャプション・注釈
FS_FOOTER = 9              # フッター（フッターのみ9pt許容）
FS_KPI = 48                # KPI大数字
FS_KPI_LABEL = 14          # KPIラベル

# 互換エイリアス（旧名→新名）
FS_SLIDE_TITLE = 14        # ヘッダー内トピック名（小さいラベル）
FS_MESSAGE = FS_ACTION_TITLE  # リード文（So What）
FS_SECTION_HEADING = FS_HEADING
FS_SUB = FS_BODY_SMALL
FS_LABEL = 11              # 最低可読サイズ（投影環境対応）

WHITE = RGBColor(255, 255, 255)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 色取得ユーティリティ
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _c(design, key):
    """DesignConfigからRGBColor を返す。"""
    return RGBColor(*design.get_color_rgb(key))


def _fade_color(design, index, total, base_key="primary",
                fade_target=220, fade_strength=0.6):
    """インデックスに応じて色をフェードさせる。上位ほど濃い。"""
    rgb = design.get_color_rgb(base_key)
    fade = index / max(total - 1, 1)
    r = int(rgb[0] + (fade_target - rgb[0]) * fade * fade_strength)
    g = int(rgb[1] + (fade_target - rgb[1]) * fade * fade_strength)
    b = int(rgb[2] + (fade_target - rgb[2]) * fade * fade_strength)
    return RGBColor(min(r, 255), min(g, 255), min(b, 255))


def _safe(val, minimum=0.05):
    """寸法を最小正数でクランプする（負やゼロを防止）。"""
    return max(val, minimum)


def _truncate(text, max_chars):
    """テキストを最大文字数で切り詰める。"""
    if not text or len(text) <= max_chars:
        return text
    return text[:max_chars - 1] + "…"


def _estimate_lines(text, font_size_pt, box_width_inches):
    """テキストが何行に折り返されるかを推定する。

    日本語: 1文字 ≈ font_size/72 inches
    英語: 1文字 ≈ font_size/72 * 0.6 inches
    混合テキストは0.8倍で概算する。
    """
    if not text:
        return 0
    char_w = font_size_pt / 72 * 0.8  # 混合推定
    chars_per_line = max(int(box_width_inches / char_w), 1)
    return max(1, -(-len(text) // chars_per_line))  # ceiling division


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 描画プリミティブ
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _text_box(slide, left, top, width, height, text, font_name, size, bold,
              color, align, *, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """テキストボックスを追加する。垂直アンカー・行間を制御可能。"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = Pt(int(size * line_spacing))

    # 垂直方向のアンカー — bodyPr 要素に設定する
    try:
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        body_pr = tf._txBody.find('.//a:bodyPr', nsmap)
        if body_pr is not None:
            anchor_val = {
                MSO_ANCHOR.TOP: 't',
                MSO_ANCHOR.MIDDLE: 'ctr',
                MSO_ANCHOR.BOTTOM: 'b',
            }.get(anchor, 't')
            body_pr.set('anchor', anchor_val)
            # spAutoFit を除去し、固定サイズにする
            for auto_fit in body_pr.findall('.//a:spAutoFit', nsmap):
                body_pr.remove(auto_fit)
    except Exception:
        pass

    return txBox


def _rect(slide, left, top, width, height, fill_color, *,
          line_color=None, line_width=None):
    """矩形シェイプを追加する。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def _circle(slide, cx, cy, radius, fill_color):
    """円形シェイプを追加する。cx, cy は中心座標。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - radius), Inches(cy - radius),
        Inches(radius * 2), Inches(radius * 2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _triangle(slide, left, top, width, height, fill_color, rotation=0):
    """三角形シェイプを追加する。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if rotation:
        shape.rotation = rotation
    return shape


def _hline(slide, left, top, width, color, thickness=1.5):
    """水平線を追加する。"""
    return _rect(slide, left, top, width, thickness / 72, color)


def _shape_text(shape, text, size, bold, color, align=PP_ALIGN.CENTER):
    """シェイプ内にテキストを設定する。"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Noto Sans JP"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def _style_cell(cell, size, bold, color):
    """テーブルセルのフォントスタイルを設定する。"""
    for p in cell.text_frame.paragraphs:
        p.font.name = "Noto Sans JP"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)


def _accent_line(slide, left, top, width, design):
    """アクセントカラーの太線。"""
    return _hline(slide, left, top, width, _c(design, "highlight"), 3)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 共通パーツ: ヘッダー / フッター / アクセント / ノート
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _add_header(slide, section, design):
    """コンサル標準ヘッダー: アクションタイトル + アクセントライン。

    spec §3-3, §5-1 準拠:
      Action Title: top=0.30, height=0.83, 20pt Bold, Primary色
      Accent Line:  top=1.15, width=12.333", アクセントカラー, 1.2pt
    """
    slide_title = section.get("slide_title", "")
    message = section.get("title", "")

    # アクションタイトル（So What主張文）— spec §5-1: 20pt Bold, 最大2行
    # slide_title がある場合は小さいトピック名 + 大きい主張文の2段構成
    if slide_title and message:
        # 上段: トピック名（小さめ、Primary色）
        _text_box(slide, MARGIN_L, HEADER_TOP, CONTENT_W, 0.30,
                  slide_title, "Noto Sans JP", FS_SLIDE_TITLE, True,
                  _c(design, "primary"), PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP)
        # 下段: 主張文（アクションタイトル）
        _text_box(slide, MARGIN_L, HEADER_TOP + 0.30, CONTENT_W, HEADER_HEIGHT - 0.30,
                  message, "Noto Sans JP", FS_ACTION_TITLE, True,
                  _c(design, "text"), PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.2)
    elif message:
        # 主張文のみ（フル高さ）
        _text_box(slide, MARGIN_L, HEADER_TOP, CONTENT_W, HEADER_HEIGHT,
                  message, "Noto Sans JP", FS_ACTION_TITLE, True,
                  _c(design, "text"), PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    elif slide_title:
        _text_box(slide, MARGIN_L, HEADER_TOP, CONTENT_W, HEADER_HEIGHT,
                  slide_title, "Noto Sans JP", FS_ACTION_TITLE, True,
                  _c(design, "text"), PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.MIDDLE)

    # アクセント区切り線 — spec §5-2: 1.2pt, アクセントカラー
    _hline(slide, MARGIN_L, ACCENT_LINE_TOP, CONTENT_W,
           _c(design, "highlight"), 1.2)


def _add_footer(slide, source_text, design, page_num=None):
    """コンサル標準フッター: 3セクション構成。

    spec §5-3 準拠:
      左: 機密表示 "CONFIDENTIAL"
      中: 出典 / プロジェクト名
      右: ページ番号（右揃え）
      全て 9pt, #999999
    """
    footer_color = _c(design, "light_text")

    # 上部区切り線（薄いボーダー）
    _hline(slide, MARGIN_L, FOOTER_TOP - 0.02, CONTENT_W,
           _c(design, "border"), 0.5)

    # 左: 機密表示
    _text_box(slide, MARGIN_L, FOOTER_TOP, 4.0, FOOTER_H,
              "CONFIDENTIAL", "Noto Sans JP", FS_FOOTER, False,
              footer_color, PP_ALIGN.LEFT)

    # 中央: 出典
    if source_text:
        _text_box(slide, MARGIN_L + 4.0, FOOTER_TOP,
                  CONTENT_W - 6.0, FOOTER_H,
                  f"Source: {source_text}", "Noto Sans JP", FS_FOOTER, False,
                  footer_color, PP_ALIGN.CENTER)

    # 右: ページ番号
    if page_num is not None:
        _text_box(slide, MARGIN_L + CONTENT_W - 2.0, FOOTER_TOP,
                  2.0, FOOTER_H,
                  str(page_num), "Noto Sans JP", FS_FOOTER, False,
                  footer_color, PP_ALIGN.RIGHT)


def _add_left_accent(slide, design):
    """ボディ左端の細いアクセント線。"""
    _rect(slide, 0, HEADER_BOTTOM, 3 / 72, SLIDE_H - HEADER_BOTTOM,
          _c(design, "highlight"))


def _add_notes(slide, section):
    """スピーカーノートを追加。"""
    if section.get("notes"):
        slide.notes_slide.notes_text_frame.text = section["notes"]


def _blank_slide(prs):
    """空白レイアウトのスライドを追加。"""
    return prs.slides.add_slide(prs.slide_layouts[6])
